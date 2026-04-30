"""Step 0.6 — Node bridge end-to-end integration test.

Marked ``slow`` so it doesn't run on every PR. Triggers via:

    pytest tests/integration/test_pptxgen_node_bridge.py -m slow

Requires Node + the ``pptxgenjs`` npm package available in the
environment (``check_pptxgen_available`` returns True). Skipped
gracefully when missing.

Verifies:
- ``run_pptxgen_executor`` accepts a SlideCommand JSON payload and
  returns valid .pptx bytes (PK\\x03\\x04 magic header + parseable
  via python-pptx).
- ``PptxGenJSBlockRenderer`` end-to-end with a real outline produces
  the expected slide count (cover + TOC + KPI + dividers + content +
  summary + thanks).
- Native chart commands round-trip through the bridge into editable
  charts (chart shape present in the OOXML).
"""
from __future__ import annotations

import io
import zipfile

import pytest
from pptx import Presentation

from backend.tools.report._block_renderer import render_outline
from backend.tools.report._kpi_extractor import KPIItem
from backend.tools.report._outline import (
    ChartAsset,
    ChartBlock,
    GrowthIndicatorsBlock,
    OutlineSection,
    ParagraphBlock,
    ReportOutline,
    StatsAsset,
    TableBlock,
    reset_id_counters,
)
from backend.tools.report._pptxgen_builder import check_pptxgen_available
from backend.tools.report._pptxgen_commands import (
    AddChart,
    AddText,
    NewSlide,
    serialize_commands,
)
from backend.tools.report._pptxgen_runtime import run_pptxgen_executor
from backend.tools.report._renderers.pptxgen import PptxGenJSBlockRenderer

pytestmark = [
    pytest.mark.slow,
    pytest.mark.skipif(
        not check_pptxgen_available(),
        reason="Node.js + pptxgenjs not available; install via "
               "`npm install -g pptxgenjs` to run.",
    ),
]


_PPTX_MAGIC = b"PK\x03\x04"


# ---------------------------------------------------------------------------
# run_pptxgen_executor smoke tests
# ---------------------------------------------------------------------------

def test_executor_handles_minimal_command_stream():
    """Smallest valid program: cover slide with one text box."""
    commands = [
        NewSlide(background="1E3A5F"),
        AddText(
            x=1, y=2, w=8, h=1.5, text="集成测试",
            font_size=44, bold=True, color="FFFFFF", alignment="center",
        ),
    ]
    payload = serialize_commands(commands)
    data = run_pptxgen_executor(payload)
    assert data[:4] == _PPTX_MAGIC
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 1


def test_executor_renders_native_chart():
    """AddChart command should produce an embedded OOXML chart part."""
    commands = [
        NewSlide(),
        AddChart(
            x=1, y=1, w=8, h=5,
            chart_type="BAR",
            data=[{
                "name": "吞吐量",
                "labels": ["大连港", "营口港", "锦州港"],
                "values": [4500.5, 3200.1, 1800.0],
            }],
            options={
                "chartColors": ["1E3A5F"],
                "showLegend": False,
                "showTitle": False,
                "barDir": "col",
            },
        ),
    ]
    payload = serialize_commands(commands)
    data = run_pptxgen_executor(payload)

    # Inspect the OOXML zip to confirm a chart part was embedded.
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
    chart_parts = [n for n in names if n.startswith("ppt/charts/chart")]
    assert chart_parts, (
        f"No chart part found in OOXML. Available: {names[:20]}"
    )


def test_executor_rejects_non_array_payload():
    with pytest.raises(RuntimeError, match="exited"):
        run_pptxgen_executor('{"not": "an array"}')


def test_executor_rejects_unknown_command_type():
    payload = '[{"type": "magic", "x": 0}]'
    with pytest.raises(RuntimeError, match="Unknown command"):
        run_pptxgen_executor(payload)


# ---------------------------------------------------------------------------
# Renderer end-to-end
# ---------------------------------------------------------------------------

def _normal_outline() -> ReportOutline:
    """Mirrors test_pptxgen_block_renderer fixture; integration scope
    keeps it self-contained so this test can run independently."""
    reset_id_counters()
    chart_asset = ChartAsset(
        asset_id="C0001", source_task="T002",
        option={
            "title": {"text": "港区吞吐量"},
            "xAxis": {"type": "category",
                      "data": ["大连港", "营口港", "锦州港"]},
            "yAxis": {"type": "value"},
            "series": [{"type": "bar", "data": [4500.5, 3200.1, 1800.0]}],
        },
    )
    stats_asset = StatsAsset(
        asset_id="S0001", source_task="T003",
        summary_stats={
            "throughput": {"max": 4500.5, "min": 1800.0, "mean": 3166.87},
        },
    )
    return ReportOutline(
        metadata={
            "title": "Integration Test Report",
            "author": "CI",
            "date": "2026-04-29",
            "intent": "integration",
        },
        kpi_summary=[
            KPIItem(label="总吞吐量", value="9500.6 万吨",
                    sub="2026 Q1", trend="positive"),
            KPIItem(label="同比增长", value="12.0%",
                    sub="YoY", trend="positive"),
        ],
        sections=[
            OutlineSection(
                name="一、现状", role="status",
                blocks=[ChartBlock(block_id="B1", asset_id="C0001")],
            ),
            OutlineSection(
                name="二、分析", role="status",
                blocks=[
                    ParagraphBlock(block_id="B2", text="吞吐量稳健增长。"),
                    TableBlock(block_id="B3", asset_id="S0001"),
                    GrowthIndicatorsBlock(
                        block_id="B4",
                        growth_rates={"throughput": {"yoy": 0.12, "mom": 0.03}},
                    ),
                ],
            ),
            OutlineSection(
                name="总结与建议", role="appendix",
                blocks=[
                    ParagraphBlock(
                        block_id="B5",
                        text="三港区集装箱吞吐量稳健增长。",
                        style="lead",
                    ),
                ],
            ),
        ],
        assets={"C0001": chart_asset, "S0001": stats_asset},
    )


def test_renderer_produces_valid_pptx_end_to_end():
    """Full pipeline: outline → SlideCommands → Node bridge → .pptx bytes."""
    outline = _normal_outline()
    renderer = PptxGenJSBlockRenderer()
    data = render_outline(outline, renderer)

    assert isinstance(data, bytes)
    assert data[:4] == _PPTX_MAGIC
    assert len(data) > 50_000, "PPTX with native chart should be > 50KB"


def test_renderer_slide_count_matches_section_composition():
    """Expected: cover + TOC + KPI overview + 2 dividers + 2 content
    sets + summary + thanks ≥ 8 slides."""
    outline = _normal_outline()
    renderer = PptxGenJSBlockRenderer()
    data = render_outline(outline, renderer)
    prs = Presentation(io.BytesIO(data))
    # cover, toc, kpi overview, 2× divider, 2× content cluster, summary, thanks
    assert len(prs.slides) >= 8


def test_renderer_emits_native_chart_part():
    outline = _normal_outline()
    renderer = PptxGenJSBlockRenderer()
    data = render_outline(outline, renderer)
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
    chart_parts = [n for n in names if n.startswith("ppt/charts/chart")]
    assert chart_parts, (
        "Expected at least one ppt/charts/chartN.xml — chart was rendered "
        "as image / placeholder instead of native editable chart."
    )
