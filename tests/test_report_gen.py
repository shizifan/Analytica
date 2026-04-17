"""Report generation skills — end-to-end test.

Constructs a rich mock analysis context (simulating a complete data-analysis
pipeline with throughput statistics, growth rates, bar/line charts, raw
DataFrames, and LLM-generated summaries), then generates DOCX / PPTX / HTML
reports and validates each output structurally.

Run:
    pytest tests/test_report_gen.py -v
"""
from __future__ import annotations

import io
import re
from typing import Any

import pandas as pd
import pytest

from backend.skills.base import SkillInput, SkillOutput
from backend.skills.report._content_collector import (
    collect_and_associate,
    ReportContent, SectionContent,
    NarrativeItem, StatsTableItem, GrowthItem,
    ChartDataItem, DataFrameItem, SummaryTextItem,
)

# ═══════════════════════════════════════════════════════════════════════════
# Mock data — simulates a "2024年港口年度运营分析" pipeline
# ═══════════════════════════════════════════════════════════════════════════

class _MockOutput:
    """Lightweight stand-in for SkillOutput used in context dict."""
    def __init__(self, data: Any):
        self.data = data


# ── 1. 吞吐量描述性分析（skill_desc_analysis）──────────────────────

THROUGHPUT_NARRATIVE = (
    "2024年港口货物吞吐量整体呈稳健增长态势。全年累计完成货物吞吐量14,580万吨，"
    "同比增长8.3%，显著高于行业平均水平5.1%。\n\n"
    "分季度来看，一季度受春节假期及寒潮天气影响，吞吐量为3,200万吨，环比下降12.5%；"
    "二季度随着复工复产全面推进，吞吐量迅速恢复至3,650万吨；三季度在暑期运输旺季拉动下，"
    "单月峰值达到1,350万吨（9月），创历史新高；四季度受全球航运市场景气回暖影响，"
    "维持在3,800万吨的高位运行。\n\n"
    "从货种结构看，集装箱吞吐量增速最为亮眼，累计完成98.5万TEU，同比增长12.1%，"
    "主要得益于新航线开辟和腹地经济活跃。散杂货受国际铁矿石价格波动影响，"
    "增速相对温和（+3.7%），但仍实现正增长。液体散货因油品进口量增加，同比增长6.8%。"
)

THROUGHPUT_STATS = {
    "货物吞吐量(万吨)": {
        "mean": 1215.0, "median": 1195.0, "std": 145.3,
        "min": 980.0, "max": 1350.0, "missing_rate": 0.0,
    },
    "集装箱(万TEU)": {
        "mean": 8.21, "median": 8.05, "std": 1.12,
        "min": 6.8, "max": 10.2, "missing_rate": 0.0,
    },
    "散杂货(万吨)": {
        "mean": 420.5, "median": 415.0, "std": 38.7,
        "min": 365.0, "max": 498.0, "missing_rate": 0.0,
    },
    "液体散货(万吨)": {
        "mean": 310.2, "median": 305.0, "std": 28.4,
        "min": 270.0, "max": 368.0, "missing_rate": 0.0,
    },
}

THROUGHPUT_GROWTH = {
    "货物吞吐量(万吨)": {"yoy": 0.083, "mom": 0.025},
    "集装箱(万TEU)":    {"yoy": 0.121, "mom": 0.038},
    "散杂货(万吨)":     {"yoy": 0.037, "mom": -0.012},
    "液体散货(万吨)":   {"yoy": 0.068, "mom": 0.015},
}

# ── 2. 泊位利用率描述性分析 ─────────────────────────────────────

BERTH_NARRATIVE = (
    "泊位综合利用率方面，全年平均泊位利用率为72.8%，较上年同期提升3.2个百分点，"
    "显示港口资源调配能力持续优化。\n\n"
    "各港区表现分化显著：大连港区凭借深水泊位优势和智慧调度系统，利用率高达81.3%，"
    "居各港区首位；营口港区因鲅鱼圈港区扩建后产能释放，利用率为73.5%，"
    "同比上升5.1个百分点；锦州港区受部分泊位升级改造影响，利用率为62.4%，"
    "短期下滑但年末已恢复至68%水平。\n\n"
    "值得关注的是，高峰时段（8-10月）泊位利用率峰值达92.1%，"
    "已接近国际港口85-95%的拥堵预警线，建议加快泊位智能化改造和船舶调度算法升级。"
)

BERTH_STATS = {
    "大连港区利用率": {
        "mean": 0.813, "median": 0.825, "std": 0.065,
        "min": 0.680, "max": 0.921, "missing_rate": 0.0,
    },
    "营口港区利用率": {
        "mean": 0.735, "median": 0.742, "std": 0.058,
        "min": 0.620, "max": 0.845, "missing_rate": 0.0,
    },
    "锦州港区利用率": {
        "mean": 0.624, "median": 0.618, "std": 0.072,
        "min": 0.510, "max": 0.780, "missing_rate": 0.0,
    },
}

BERTH_GROWTH = {
    "大连港区利用率": {"yoy": 0.041, "mom": 0.018},
    "营口港区利用率": {"yoy": 0.051, "mom": -0.008},
    "锦州港区利用率": {"yoy": -0.023, "mom": 0.032},
}

# ── 3. ECharts 柱状图 — 月度吞吐量对比（skill_chart_bar）────────

MONTHLY_BAR_CHART = {
    "title": {"text": "月度货物吞吐量同比对比", "left": "center",
              "textStyle": {"color": "#1E3A5F"}},
    "tooltip": {"trigger": "axis", "axisPointer": {"type": "shadow"}},
    "legend": {"data": ["2023年", "2024年"], "bottom": 0},
    "xAxis": {
        "type": "category",
        "data": ["1月", "2月", "3月", "4月", "5月", "6月",
                 "7月", "8月", "9月", "10月", "11月", "12月"],
    },
    "yAxis": {"type": "value", "name": "万吨"},
    "series": [
        {
            "name": "2023年", "type": "bar",
            "data": [980, 920, 1050, 1080, 1100, 1120,
                     1150, 1200, 1250, 1180, 1130, 1100],
            "itemStyle": {"color": "#1E3A5F"},
        },
        {
            "name": "2024年", "type": "bar",
            "data": [1050, 980, 1130, 1170, 1200, 1230,
                     1260, 1300, 1350, 1280, 1230, 1200],
            "itemStyle": {"color": "#F0A500"},
        },
    ],
}

# ── 4. ECharts 折线图 — 泊位利用率趋势（skill_chart_line）───────

BERTH_LINE_CHART = {
    "title": {"text": "各港区泊位利用率月度趋势", "left": "center"},
    "tooltip": {"trigger": "axis"},
    "legend": {"data": ["大连港区", "营口港区", "锦州港区"], "bottom": 0},
    "xAxis": {
        "type": "category",
        "data": ["1月", "2月", "3月", "4月", "5月", "6月",
                 "7月", "8月", "9月", "10月", "11月", "12月"],
    },
    "yAxis": {"type": "value", "name": "%", "max": 100},
    "series": [
        {"name": "大连港区", "type": "line",
         "data": [72, 68, 75, 78, 80, 83, 85, 88, 92, 86, 82, 78]},
        {"name": "营口港区", "type": "line",
         "data": [65, 62, 68, 72, 74, 76, 78, 82, 85, 78, 74, 70]},
        {"name": "锦州港区", "type": "line",
         "data": [58, 51, 55, 60, 63, 65, 68, 72, 78, 70, 65, 62]},
    ],
}

# ── 5. 原始数据 DataFrame（skill_api_fetch）──────────────────────

RAW_THROUGHPUT_DF = pd.DataFrame({
    "月份":       ["2024-01", "2024-02", "2024-03", "2024-04", "2024-05", "2024-06",
                   "2024-07", "2024-08", "2024-09", "2024-10", "2024-11", "2024-12"],
    "货物吞吐量": [1050, 980, 1130, 1170, 1200, 1230, 1260, 1300, 1350, 1280, 1230, 1200],
    "集装箱TEU":  [6.8, 6.5, 7.2, 7.8, 8.1, 8.5, 8.9, 9.3, 10.2, 9.5, 8.8, 8.2],
    "泊位利用率": [0.72, 0.68, 0.75, 0.78, 0.80, 0.83, 0.85, 0.88, 0.92, 0.86, 0.82, 0.78],
})

# ── 6. LLM 生成的总结摘要（skill_summary_gen）────────────────────

SUMMARY_TEXT = (
    "综合分析表明，2024年港口运营实现量质齐升：货物吞吐量同比增长8.3%，"
    "集装箱增速达12.1%，泊位利用率提升3.2个百分点。主要驱动因素包括："
    "（1）新航线开辟带动集装箱业务快速增长；"
    "（2）智慧调度系统上线有效提升泊位周转效率；"
    "（3）腹地经济持续向好支撑货源稳定增长。"
    "建议重点关注：高峰时段泊位利用率已逼近预警线，需加快智能化改造；"
    "锦州港区升级改造后的产能释放值得期待；散杂货业务需关注国际大宗商品价格波动风险。"
)


def _build_rich_context() -> dict[str, _MockOutput]:
    """Assemble a complete mock execution context."""
    return {
        # 描述性分析 — 吞吐量
        "T001_desc_throughput": _MockOutput({
            "narrative": THROUGHPUT_NARRATIVE,
            "summary_stats": THROUGHPUT_STATS,
            "growth_rates": THROUGHPUT_GROWTH,
        }),
        # 描述性分析 — 泊位利用率
        "T002_desc_berth": _MockOutput({
            "narrative": BERTH_NARRATIVE,
            "summary_stats": BERTH_STATS,
            "growth_rates": BERTH_GROWTH,
        }),
        # 柱状图 — 月度吞吐量对比
        "T003_chart_bar_monthly": _MockOutput(MONTHLY_BAR_CHART),
        # 折线图 — 泊位利用率趋势
        "T004_chart_line_berth": _MockOutput(BERTH_LINE_CHART),
        # 原始数据 DataFrame
        "T005_raw_data": _MockOutput(RAW_THROUGHPUT_DF),
        # LLM 总结摘要
        "T006_summary": _MockOutput(SUMMARY_TEXT),
    }


REPORT_PARAMS = {
    "report_metadata": {
        "title": "2024年港口年度运营分析报告",
        "author": "数据分析团队 · Analytica",
        "date": "2024-12-31",
    },
    "report_structure": {
        "sections": [
            {"name": "货物吞吐量综合分析", "task_refs": ["T001_desc_throughput"]},
            {"name": "泊位利用率与运营效率", "task_refs": ["T002_desc_berth"]},
            {"name": "月度趋势对比", "task_refs": ["T003_chart_bar_monthly", "T004_chart_line_berth"]},
            {"name": "原始数据明细", "task_refs": ["T005_raw_data"]},
        ],
    },
}


# ═══════════════════════════════════════════════════════════════════════════
# Fixture
# ═══════════════════════════════════════════════════════════════════════════

@pytest.fixture
def rich_context():
    return _build_rich_context()


@pytest.fixture
def skill_input():
    return SkillInput(params=REPORT_PARAMS)


# ═══════════════════════════════════════════════════════════════════════════
# 1. Content collector 单元测试
# ═══════════════════════════════════════════════════════════════════════════

class TestContentCollector:
    """Verify that _content_collector correctly classifies and associates."""

    def test_extract_all_content_types(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)

        assert isinstance(report, ReportContent)
        assert report.title == "2024年港口年度运营分析报告"
        assert report.author == "数据分析团队 · Analytica"
        assert report.date == "2024-12-31"
        assert len(report.sections) == 4

    def test_every_section_has_items(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        for sec in report.sections:
            assert len(sec.items) >= 1, f"章节 '{sec.name}' 无内容项"

    def test_narrative_items_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        all_items = [it for sec in report.sections for it in sec.items]
        narratives = [it for it in all_items if isinstance(it, NarrativeItem)]
        assert len(narratives) >= 2, "应至少包含吞吐量和泊位两个 narrative"

    def test_stats_items_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        all_items = [it for sec in report.sections for it in sec.items]
        stats = [it for it in all_items if isinstance(it, StatsTableItem)]
        assert len(stats) >= 2, "应至少包含吞吐量和泊位两组统计数据"

    def test_growth_items_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        all_items = [it for sec in report.sections for it in sec.items]
        growths = [it for it in all_items if isinstance(it, GrowthItem)]
        assert len(growths) >= 2

    def test_chart_items_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        all_items = [it for sec in report.sections for it in sec.items]
        charts = [it for it in all_items if isinstance(it, ChartDataItem)]
        assert len(charts) >= 2, "应包含柱状图和折线图"

    def test_dataframe_item_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        all_items = [it for sec in report.sections for it in sec.items]
        dfs = [it for it in all_items if isinstance(it, DataFrameItem)]
        assert len(dfs) >= 1, "应包含原始 DataFrame"

    def test_summary_items_present(self, rich_context):
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        assert len(report.summary_items) >= 1, "应包含 LLM 总结摘要"

    def test_task_refs_throughput_section(self, rich_context):
        """吞吐量 task 的内容应通过 task_refs 精确映射到第 1 章节。"""
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        tp_section = report.sections[0]
        assert tp_section.name == "货物吞吐量综合分析"
        source_tasks = {it.source_task for it in tp_section.items}
        assert "T001_desc_throughput" in source_tasks
        item_types = {type(it).__name__ for it in tp_section.items}
        assert "NarrativeItem" in item_types
        assert "StatsTableItem" in item_types
        assert "GrowthItem" in item_types

    def test_task_refs_berth_section(self, rich_context):
        """泊位 task 的内容应通过 task_refs 精确映射到第 2 章节。"""
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        berth_section = report.sections[1]
        assert berth_section.name == "泊位利用率与运营效率"
        source_tasks = {it.source_task for it in berth_section.items}
        assert "T002_desc_berth" in source_tasks
        item_types = {type(it).__name__ for it in berth_section.items}
        assert "NarrativeItem" in item_types
        assert "StatsTableItem" in item_types

    def test_task_refs_chart_section(self, rich_context):
        """图表 task 应通过 task_refs 映射到第 3 章节（月度趋势对比）。"""
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        chart_section = report.sections[2]
        assert chart_section.name == "月度趋势对比"
        source_tasks = {it.source_task for it in chart_section.items}
        assert "T003_chart_bar_monthly" in source_tasks
        assert "T004_chart_line_berth" in source_tasks
        item_types = {type(it).__name__ for it in chart_section.items}
        assert "ChartDataItem" in item_types

    def test_task_refs_dataframe_section(self, rich_context):
        """DataFrame task 应通过 task_refs 映射到第 4 章节（原始数据明细）。"""
        report = collect_and_associate(REPORT_PARAMS, rich_context)
        df_section = report.sections[3]
        assert df_section.name == "原始数据明细"
        source_tasks = {it.source_task for it in df_section.items}
        assert "T005_raw_data" in source_tasks
        item_types = {type(it).__name__ for it in df_section.items}
        assert "DataFrameItem" in item_types

    def test_fallback_sequential_without_task_refs(self, rich_context):
        """当 sections 为旧格式（纯字符串列表）时，应按 task 顺序分配。"""
        params = {
            "report_metadata": {"title": "Fallback 测试"},
            "report_structure": {
                "sections": ["章节A", "章节B", "章节C"],
            },
        }
        report = collect_and_associate(params, rich_context)
        assert len(report.sections) == 3
        # 所有内容都应被分配（无遗漏）
        total_items = sum(len(sec.items) for sec in report.sections)
        assert total_items >= 5, f"Fallback 模式下应分配所有内容项，实际 {total_items}"

    def test_auto_sections_when_empty(self, rich_context):
        """当 sections 为空时，应根据 context 自动生成。"""
        params = {
            "report_metadata": {"title": "自动生成章节测试"},
            "report_structure": {"sections": []},
        }
        report = collect_and_associate(params, rich_context)
        assert len(report.sections) >= 1


# ═══════════════════════════════════════════════════════════════════════════
# 2. DOCX 生成测试
# ═══════════════════════════════════════════════════════════════════════════

class TestDocxGeneration:
    """End-to-end DOCX generation and structural validation."""

    async def test_docx_generates_successfully(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)

        assert result.status == "success"
        assert result.output_type == "file"
        assert isinstance(result.data, bytes)
        assert result.metadata["format"] == "docx"
        assert result.metadata["file_size_bytes"] > 0

    async def test_docx_parseable(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)

        doc = Document(io.BytesIO(result.data))
        assert len(doc.paragraphs) > 0

    async def test_docx_has_headings(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)
        doc = Document(io.BytesIO(result.data))

        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert any("货物吞吐量" in h for h in headings), f"缺少吞吐量标题, got: {headings}"
        assert any("泊位利用率" in h for h in headings), f"缺少泊位标题, got: {headings}"
        assert any("总结" in h for h in headings), f"缺少总结标题, got: {headings}"

    async def test_docx_has_tables(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)
        doc = Document(io.BytesIO(result.data))

        # 应至少有: 吞吐量统计表、泊位统计表、月度图表数据表、DataFrame 表
        assert len(doc.tables) >= 3, f"表格数量不足: {len(doc.tables)}"

    async def test_docx_contains_narrative_text(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)
        doc = Document(io.BytesIO(result.data))

        full_text = "\n".join(p.text for p in doc.paragraphs)
        assert "14,580万吨" in full_text or "14580" in full_text, "narrative 内容未写入"

    async def test_docx_cover_page(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)
        doc = Document(io.BytesIO(result.data))

        full_text = "\n".join(p.text for p in doc.paragraphs[:10])
        assert "2024年港口年度运营分析报告" in full_text
        assert "数据分析团队" in full_text

    async def test_docx_has_page_header(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from docx import Document

        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(skill_input, rich_context)
        doc = Document(io.BytesIO(result.data))

        section = doc.sections[0]
        header_text = section.header.paragraphs[0].text if section.header.paragraphs else ""
        assert "2024年港口年度运营分析报告" in header_text


# ═══════════════════════════════════════════════════════════════════════════
# 3. PPTX 生成测试
# ═══════════════════════════════════════════════════════════════════════════

class TestPptxGeneration:
    """End-to-end PPTX generation and structural validation."""

    async def test_pptx_generates_successfully(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)

        assert result.status == "success"
        assert result.output_type == "file"
        assert isinstance(result.data, bytes)
        assert result.metadata["format"] == "pptx"
        assert result.metadata["slide_count"] > 0

    async def test_pptx_parseable(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)

        prs = Presentation(io.BytesIO(result.data))
        assert len(prs.slides) > 0

    async def test_pptx_slide_count_reasonable(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)

        count = result.metadata["slide_count"]
        # Cover + TOC + (divider+content)*4_sections + summary + thankyou
        assert count >= 10, f"幻灯片太少: {count}"
        assert count <= 40, f"幻灯片太多: {count}"

    async def test_pptx_has_cover_slide(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        slide1_texts = _pptx_slide_texts(prs.slides[0])
        assert any("2024年港口年度运营分析报告" in t for t in slide1_texts)

    async def test_pptx_has_toc_slide(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        slide2_texts = _pptx_slide_texts(prs.slides[1])
        assert any("目" in t for t in slide2_texts), "第 2 张应为目录页"

    async def test_pptx_has_section_divider(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        slide3_texts = _pptx_slide_texts(prs.slides[2])
        assert any("01" in t for t in slide3_texts), "第 3 张应为章节分隔页"

    async def test_pptx_has_summary_slide(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        slides_list = list(prs.slides)
        last_two = [_pptx_slide_texts(s) for s in slides_list[-2:]]
        all_last_texts = [t for texts in last_two for t in texts]
        assert any("核心结论" in t or "总结" in t for t in all_last_texts), \
            f"倒数第 2 张应为总结页, got: {all_last_texts}"

    async def test_pptx_has_thankyou_slide(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        last_texts = _pptx_slide_texts(list(prs.slides)[-1])
        assert any("谢谢" in t or "THANK" in t for t in last_texts)

    async def test_pptx_has_tables(self, skill_input, rich_context):
        from backend.skills.report.pptx_gen import PptxReportSkill
        from pptx import Presentation

        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(skill_input, rich_context)
        prs = Presentation(io.BytesIO(result.data))

        table_count = sum(
            1 for slide in prs.slides for shape in slide.shapes if shape.has_table
        )
        assert table_count >= 2, f"PPTX 中表格太少: {table_count}"


# ═══════════════════════════════════════════════════════════════════════════
# 4. HTML 生成测试
# ═══════════════════════════════════════════════════════════════════════════

class TestHtmlGeneration:
    """End-to-end HTML generation and content validation."""

    async def test_html_generates_successfully(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert result.status == "success"
        assert result.output_type == "file"
        assert isinstance(result.data, str)
        assert result.metadata["format"] == "html"

    async def test_html_is_valid_structure(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)
        html = result.data

        assert "<!DOCTYPE html>" in html
        assert "<html" in html
        assert "</html>" in html
        assert "<title>2024年港口年度运营分析报告</title>" in html

    async def test_html_has_sections(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)
        html = result.data

        assert "货物吞吐量综合分析" in html
        assert "泊位利用率与运营效率" in html
        assert "月度趋势对比" in html
        assert "原始数据明细" in html

    async def test_html_has_narrative(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert "14,580万吨" in result.data or "14580" in result.data

    async def test_html_has_echarts(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert "echarts.init" in result.data
        assert result.metadata["chart_count"] >= 2, "应包含至少 2 个 ECharts 图表"

    async def test_html_has_stats_table(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert 'class="stats"' in result.data, "应包含统计表格"
        assert "mean" in result.data

    async def test_html_has_kpi_cards(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert "kpi-card" in result.data, "应包含 KPI 数据卡片"
        assert "kpi-row" in result.data

    async def test_html_has_summary(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert "总结与建议" in result.data
        assert "智慧调度" in result.data or "智能化" in result.data

    async def test_html_has_dataframe_table(self, skill_input, rich_context):
        from backend.skills.report.html_gen import HtmlReportSkill

        skill = HtmlReportSkill()
        skill.skill_id = "skill_report_html"
        result = await skill.execute(skill_input, rich_context)

        assert "数据明细" in result.data
        assert "2024-01" in result.data


# ═══════════════════════════════════════════════════════════════════════════
# 5. 跨格式一致性测试
# ═══════════════════════════════════════════════════════════════════════════

class TestCrossFormatConsistency:
    """Verify that all three formats receive the same structured content."""

    async def test_all_three_formats_succeed(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from backend.skills.report.pptx_gen import PptxReportSkill
        from backend.skills.report.html_gen import HtmlReportSkill

        results = {}
        for SkillCls, name in [
            (DocxReportSkill, "docx"),
            (PptxReportSkill, "pptx"),
            (HtmlReportSkill, "html"),
        ]:
            skill = SkillCls()
            skill.skill_id = f"skill_report_{name}"
            results[name] = await skill.execute(skill_input, rich_context)

        for name, r in results.items():
            assert r.status == "success", f"{name} 生成失败: {r.error_message}"

    async def test_file_sizes_reasonable(self, skill_input, rich_context):
        from backend.skills.report.docx_gen import DocxReportSkill
        from backend.skills.report.pptx_gen import PptxReportSkill
        from backend.skills.report.html_gen import HtmlReportSkill

        for SkillCls, name, min_size in [
            (DocxReportSkill, "docx", 10_000),
            (PptxReportSkill, "pptx", 20_000),
            (HtmlReportSkill, "html", 2_000),
        ]:
            skill = SkillCls()
            skill.skill_id = f"skill_report_{name}"
            result = await skill.execute(skill_input, rich_context)
            size = result.metadata.get("file_size_bytes", len(result.data))
            assert size >= min_size, f"{name} 文件太小: {size} bytes"


# ═══════════════════════════════════════════════════════════════════════════
# 6. 边界条件测试
# ═══════════════════════════════════════════════════════════════════════════

class TestEdgeCases:
    """Test graceful handling of minimal / empty input."""

    async def test_empty_context(self):
        """空 context 不应崩溃。"""
        from backend.skills.report.docx_gen import DocxReportSkill
        from backend.skills.report.pptx_gen import PptxReportSkill
        from backend.skills.report.html_gen import HtmlReportSkill

        inp = SkillInput(params={
            "report_metadata": {"title": "空报告测试"},
            "report_structure": {"sections": ["概览"]},
        })
        empty_ctx: dict = {}

        for SkillCls, name in [
            (DocxReportSkill, "docx"),
            (PptxReportSkill, "pptx"),
            (HtmlReportSkill, "html"),
        ]:
            skill = SkillCls()
            skill.skill_id = f"skill_report_{name}"
            result = await skill.execute(inp, empty_ctx)
            assert result.status == "success", f"{name} 空 context 失败: {result.error_message}"

    async def test_no_sections_provided(self, rich_context):
        """未提供 sections 时应自动生成。"""
        from backend.skills.report.docx_gen import DocxReportSkill

        inp = SkillInput(params={
            "report_metadata": {"title": "无章节测试"},
            "report_structure": {},
        })
        skill = DocxReportSkill()
        skill.skill_id = "skill_report_docx"
        result = await skill.execute(inp, rich_context)
        assert result.status == "success"

    async def test_single_narrative_only(self):
        """仅有一条 narrative 的极简场景。"""
        from backend.skills.report.pptx_gen import PptxReportSkill

        ctx = {"T001": _MockOutput({"narrative": "这是一段简短的分析文本，用于测试极简输入场景的报告生成能力。"})}
        inp = SkillInput(params={
            "report_metadata": {"title": "极简报告"},
            "report_structure": {"sections": ["分析"]},
        })
        skill = PptxReportSkill()
        skill.skill_id = "skill_report_pptx"
        result = await skill.execute(inp, ctx)
        assert result.status == "success"
        assert result.metadata["slide_count"] >= 4  # cover + toc + divider+content + summary + thankyou


# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════

def _pptx_slide_texts(slide) -> list[str]:
    """Extract all text from a PPTX slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return texts
