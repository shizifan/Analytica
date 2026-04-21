"""PPTX Report Generation Skill — Skill mode (LLM + tools) with deterministic fallback.

Orchestrator: delegates metadata/content extraction to ``_content_collector``,
then either runs an LLM agent loop that composes the presentation using tools
(wrapping ``_pptx_slides``), or falls back to a deterministic builder.
"""
from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from backend.skills.base import BaseSkill, SkillCategory, SkillInput, SkillOutput
from backend.skills.registry import register_skill
from backend.skills.report._content_collector import (
    collect_and_associate,
    NarrativeItem, StatsTableItem, GrowthItem,
    ChartDataItem, DataFrameItem, ReportContent,
)
from backend.skills.report import _pptx_slides as S
from backend.skills.report import _theme as T

logger = logging.getLogger("analytica.skills.report_pptx")


# ---------------------------------------------------------------------------
# Deterministic builder (extracted from the original execute body)
# ---------------------------------------------------------------------------

def _stats_to_text(summary_stats: dict[str, Any]) -> str:
    """Convert summary_stats dict to readable text for two-column layout."""
    lines: list[str] = []
    for col, vals in summary_stats.items():
        if not isinstance(vals, dict):
            continue
        mean = vals.get("mean")
        std = vals.get("std")
        if mean is not None:
            line = f"{col}：均值 {mean:,.2f}"
            if std is not None:
                line += f"  标准差 {std:,.2f}"
            lines.append(line)
    return "\n".join(lines) if lines else "暂无统计数据"


def _build_pptx_deterministic(prs: Presentation, report: ReportContent) -> None:
    """Build PPTX content using hardcoded slide ordering — the fallback path."""
    S.build_cover_slide(prs, report.title, report.author, report.date)
    S.build_toc_slide(prs, [s.name for s in report.sections])

    for idx, section in enumerate(report.sections, 1):
        S.build_section_divider_slide(prs, idx, section.name)

        narratives = [it for it in section.items if isinstance(it, NarrativeItem)]
        stats_items = [it for it in section.items if isinstance(it, StatsTableItem)]
        growth_items = [it for it in section.items if isinstance(it, GrowthItem)]
        chart_items = [it for it in section.items if isinstance(it, ChartDataItem)]

        for gi in growth_items:
            S.build_kpi_cards_slide(prs, section.name, gi.growth_rates)

        if narratives and stats_items:
            nar_text = "\n\n".join(n.text for n in narratives)
            stats_text = _stats_to_text(stats_items[0].summary_stats)
            S.build_two_column_slide(prs, section.name, nar_text, stats_text)
            for si in stats_items:
                S.build_stats_table_slide(prs, f"{section.name} - 统计数据", si.summary_stats)
        elif narratives:
            nar_text = "\n\n".join(n.text for n in narratives)
            S.build_narrative_slide(prs, section.name, nar_text)
        elif stats_items:
            for si in stats_items:
                S.build_stats_table_slide(prs, f"{section.name} - 统计数据", si.summary_stats)

        for ci in chart_items:
            S.build_chart_table_slide(prs, ci.option)

    conclusions = []
    for si in report.summary_items:
        text = si.text
        conclusions.append(text[:120] + "..." if len(text) > 120 else text)
    if not conclusions:
        for sec in report.sections:
            for item in sec.items:
                if isinstance(item, NarrativeItem) and len(item.text) > 20:
                    conclusions.append(item.text[:100] + "...")
                    break
    if not conclusions:
        conclusions = ["数据分析完成，详见各章节内容"]

    S.build_summary_slide(prs, conclusions)
    S.build_thank_you_slide(prs)


# ---------------------------------------------------------------------------
# Skill
# ---------------------------------------------------------------------------

@register_skill("skill_report_pptx", SkillCategory.REPORT, "PPTX 报告生成（封面/目录/图表/总结）",
                input_spec="report_metadata + report_structure + 上游数据/图表引用",
                output_spec="PPTX 文件路径")
class PptxReportSkill(BaseSkill):

    async def execute(self, inp: SkillInput, context: dict[str, Any]) -> SkillOutput:
        try:
            report = collect_and_associate(
                inp.params, context,
                task_order=inp.params.get("_task_order"),
            )

            # ── Try Skill mode (LLM agent loop) ──
            mode = "deterministic_fallback"
            prs = Presentation()
            prs.slide_width = Inches(T.SLIDE_WIDTH)
            prs.slide_height = Inches(T.SLIDE_HEIGHT)

            try:
                from backend.config import get_settings

                settings = get_settings()
                if not settings.REPORT_AGENT_ENABLED:
                    raise RuntimeError("agent mode disabled by config")

                from langchain_openai import ChatOpenAI

                from backend.skills.report._agent_loop import run_report_agent, serialize_report_content
                from backend.skills.report._pptx_tools import PPTX_SYSTEM_PROMPT, make_pptx_tools

                llm = ChatOpenAI(
                    base_url=settings.QWEN_API_BASE,
                    api_key=settings.QWEN_API_KEY,
                    model=settings.QWEN_MODEL,
                    temperature=0.2,
                    request_timeout=90,
                    extra_body={"enable_thinking": False},
                )

                tools = make_pptx_tools(prs, report)
                user_message = serialize_report_content(report)
                success = await run_report_agent(llm, tools, PPTX_SYSTEM_PROMPT, user_message)

                if success:
                    mode = "llm_agent"
                else:
                    logger.warning("PPTX agent did not finalise; falling back to deterministic")
                    prs = Presentation()
                    prs.slide_width = Inches(T.SLIDE_WIDTH)
                    prs.slide_height = Inches(T.SLIDE_HEIGHT)
                    _build_pptx_deterministic(prs, report)

            except Exception as agent_err:
                logger.warning("PPTX agent loop failed (%s); falling back to deterministic", agent_err)
                prs = Presentation()
                prs.slide_width = Inches(T.SLIDE_WIDTH)
                prs.slide_height = Inches(T.SLIDE_HEIGHT)
                _build_pptx_deterministic(prs, report)
                mode = "deterministic_fallback_error"

            # ── Serialise ──
            buffer = io.BytesIO()
            prs.save(buffer)
            pptx_bytes = buffer.getvalue()

            return SkillOutput(
                skill_id=self.skill_id,
                status="success",
                output_type="file",
                data=pptx_bytes,
                metadata={
                    "format": "pptx",
                    "slide_count": len(prs.slides),
                    "title": report.title,
                    "file_size_bytes": len(pptx_bytes),
                    "mode": mode,
                },
            )

        except Exception as e:
            logger.exception("PPTX generation failed: %s", e)
            return self._fail(str(e))
