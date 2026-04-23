"""PPTX slide builder functions (python-pptx).

Each function creates one or more slides on the given ``Presentation``.
All are stateless — the orchestrator in ``pptx_gen.py`` calls them
in the desired order.
"""
from __future__ import annotations

from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from backend.tools.report import _theme as T

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(t: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


def _set_bg(slide, color: tuple[int, int, int]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_textbox(slide, left, top, width, height, text,
                 font_size=16, bold=False, color=T.RGB_WHITE,
                 alignment=PP_ALIGN.LEFT, font_name=T.FONT_CN):
    """Convenience: add a single-paragraph text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_rect(slide, left, top, width, height, fill_color: tuple[int, int, int]):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.fill.background()
    return shape


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fmt_number(v: Any) -> str:
    if v is None:
        return "-"
    if isinstance(v, float):
        if abs(v) < 1:
            return f"{v:.4f}"
        return f"{v:,.2f}"
    return str(v)

# ---------------------------------------------------------------------------
# Cover slide
# ---------------------------------------------------------------------------

def build_cover_slide(prs: Presentation, title: str, author: str, date: str) -> None:
    slide = _blank(prs)
    _set_bg(slide, T.RGB_PRIMARY)
    _add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1.5),
                 title, font_size=T.SIZE_TITLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
                 f"编制：{author}", font_size=18,
                 color=T.RGB_ACCENT, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
                 date, font_size=14, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# TOC slide
# ---------------------------------------------------------------------------

def build_toc_slide(prs: Presentation, sections: list[str]) -> None:
    slide = _blank(prs)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "目  录", font_size=24, bold=True,
                 color=T.RGB_PRIMARY, alignment=PP_ALIGN.CENTER)
    # Left accent bar
    _add_rect(slide, Inches(1.2), Inches(1.3), Inches(0.08), Inches(min(len(sections) * 0.6, 5.5)),
              T.RGB_ACCENT)
    for i, name in enumerate(sections[:10]):
        _add_textbox(slide, Inches(1.5), Inches(1.3 + i * 0.55), Inches(7), Inches(0.5),
                     f"{i + 1}.  {name}", font_size=16, color=T.RGB_PRIMARY)

# ---------------------------------------------------------------------------
# Section divider slide
# ---------------------------------------------------------------------------

def build_section_divider_slide(prs: Presentation, number: int, title: str) -> None:
    slide = _blank(prs)
    _set_bg(slide, T.RGB_PRIMARY)
    # Large section number
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(3), Inches(2),
                 f"{number:02d}", font_size=72, bold=True,
                 color=T.RGB_ACCENT, font_name=T.FONT_NUM)
    # Title
    _add_textbox(slide, Inches(1), Inches(3.8), Inches(8), Inches(1),
                 title, font_size=T.SIZE_H1, bold=True)
    # Accent underline
    _add_rect(slide, Inches(1), Inches(4.8), Inches(2), Inches(0.06), T.RGB_ACCENT)

# ---------------------------------------------------------------------------
# Narrative slide (with left color bar)
# ---------------------------------------------------------------------------

def build_narrative_slide(prs: Presentation, title: str, text: str) -> None:
    """Create one or more narrative slides; auto-paginate long text."""
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    max_lines_per_slide = 12
    chunks = [lines[i:i + max_lines_per_slide] for i in range(0, len(lines), max_lines_per_slide)]
    if not chunks:
        chunks = [["（暂无详细内容）"]]

    for ci, chunk in enumerate(chunks):
        slide = _blank(prs)
        # Left accent bar
        _add_rect(slide, Inches(0.3), Inches(0.3), Inches(0.08), Inches(6.9), T.RGB_ACCENT)

        # Title
        suffix = f" ({ci+1}/{len(chunks)})" if len(chunks) > 1 else ""
        _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(9), Inches(0.8),
                     f"{title}{suffix}", font_size=T.SIZE_H2, bold=True,
                     color=T.RGB_PRIMARY)

        # Content
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.6), Inches(5.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        for li, line in enumerate(chunk):
            p = tf.add_paragraph() if li > 0 else tf.paragraphs[0]
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = _rgb(T.RGB_TEXT_DARK)
            p.font.name = T.FONT_CN
            p.space_before = Pt(6)

# ---------------------------------------------------------------------------
# KPI cards slide
# ---------------------------------------------------------------------------

def build_kpi_cards_slide(prs: Presentation, title: str,
                          growth_rates: dict[str, dict[str, float | None]]) -> None:
    """2-4 KPI cards in a horizontal row."""
    items = [(col, rates) for col, rates in growth_rates.items()
             if isinstance(rates, dict) and (rates.get("yoy") is not None or rates.get("mom") is not None)]
    if not items:
        return

    items = items[:4]
    slide = _blank(prs)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 title, font_size=T.SIZE_H2, bold=True,
                 color=T.RGB_PRIMARY, alignment=PP_ALIGN.CENTER)

    n = len(items)
    card_w = 8.0 / n
    start_x = (10 - card_w * n) / 2

    for i, (col, rates) in enumerate(items):
        x = start_x + i * card_w
        # Card background
        _add_rect(slide, Inches(x + 0.15), Inches(1.6), Inches(card_w - 0.3), Inches(4.5),
                  T.RGB_BG_LIGHT)

        # Label
        _add_textbox(slide, Inches(x + 0.3), Inches(1.8), Inches(card_w - 0.6), Inches(0.6),
                     col, font_size=T.SIZE_KPI_LABEL, bold=True,
                     color=T.RGB_NEUTRAL, alignment=PP_ALIGN.CENTER)

        # YoY
        yoy = rates.get("yoy")
        if yoy is not None:
            arrow = "\u2191" if yoy >= 0 else "\u2193"
            color = T.RGB_POSITIVE if yoy >= 0 else T.RGB_NEGATIVE
            _add_textbox(slide, Inches(x + 0.3), Inches(2.6), Inches(card_w - 0.6), Inches(1.2),
                         f"{arrow}{abs(yoy)*100:.1f}%", font_size=T.SIZE_KPI_LARGE, bold=True,
                         color=color, alignment=PP_ALIGN.CENTER, font_name=T.FONT_NUM)
            _add_textbox(slide, Inches(x + 0.3), Inches(3.8), Inches(card_w - 0.6), Inches(0.5),
                         "同比", font_size=T.SIZE_KPI_LABEL,
                         color=T.RGB_NEUTRAL, alignment=PP_ALIGN.CENTER)

        # MoM
        mom = rates.get("mom")
        if mom is not None:
            arrow = "\u2191" if mom >= 0 else "\u2193"
            color = T.RGB_POSITIVE if mom >= 0 else T.RGB_NEGATIVE
            _add_textbox(slide, Inches(x + 0.3), Inches(4.4), Inches(card_w - 0.6), Inches(1.0),
                         f"{arrow}{abs(mom)*100:.1f}%", font_size=24, bold=True,
                         color=color, alignment=PP_ALIGN.CENTER, font_name=T.FONT_NUM)
            _add_textbox(slide, Inches(x + 0.3), Inches(5.4), Inches(card_w - 0.6), Inches(0.5),
                         "环比", font_size=T.SIZE_KPI_LABEL,
                         color=T.RGB_NEUTRAL, alignment=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Stats table slide
# ---------------------------------------------------------------------------

def build_stats_table_slide(prs: Presentation, title: str,
                            summary_stats: dict[str, Any]) -> None:
    """Render summary_stats as a table on a slide."""
    if not summary_stats:
        return

    # Flatten grouped stats
    first_val = next(iter(summary_stats.values()))
    if isinstance(first_val, dict) and not any(k in first_val for k in ("mean", "median", "std", "min", "max")):
        flat: dict[str, dict] = {}
        for gk, cols in summary_stats.items():
            if isinstance(cols, dict):
                for cn, metrics in cols.items():
                    flat[f"{gk}/{cn}"] = metrics if isinstance(metrics, dict) else {}
        summary_stats = flat if flat else summary_stats

    metrics = ["mean", "median", "std", "min", "max"]
    col_names = [c for c, v in summary_stats.items() if isinstance(v, dict)]
    if not col_names:
        return

    slide = _blank(prs)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 title, font_size=T.SIZE_H2, bold=True,
                 color=T.RGB_PRIMARY)

    n_rows = len(col_names)
    n_cols = 1 + len(metrics)
    rows = 1 + n_rows
    table_shape = slide.shapes.add_table(rows, n_cols,
                                         Inches(0.5), Inches(1.2),
                                         Inches(9), Inches(min(n_rows * 0.55 + 0.5, 5.8)))
    table = table_shape.table

    headers = ["指标"] + metrics
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(T.RGB_PRIMARY)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(T.SIZE_TABLE_HEADER)
            p.font.bold = True
            p.font.color.rgb = _rgb(T.RGB_WHITE)
            p.font.name = T.FONT_NUM

    for ri, col in enumerate(col_names):
        vals = summary_stats[col]
        table.cell(ri + 1, 0).text = str(col)
        for mi, m in enumerate(metrics):
            table.cell(ri + 1, mi + 1).text = _fmt_number(vals.get(m))
        if ri % 2 == 1:
            for ci in range(n_cols):
                table.cell(ri + 1, ci).fill.solid()
                table.cell(ri + 1, ci).fill.fore_color.rgb = _rgb(T.RGB_BG_LIGHT)
        for ci in range(n_cols):
            for p in table.cell(ri + 1, ci).text_frame.paragraphs:
                p.font.size = Pt(T.SIZE_TABLE_BODY)
                p.font.name = T.FONT_NUM if ci > 0 else T.FONT_CN

# ---------------------------------------------------------------------------
# Two-column slide (narrative + data)
# ---------------------------------------------------------------------------

def build_two_column_slide(prs: Presentation, title: str,
                           left_text: str, right_text: str) -> None:
    slide = _blank(prs)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 title, font_size=T.SIZE_H2, bold=True, color=T.RGB_PRIMARY)

    # Divider line
    _add_rect(slide, Inches(4.95), Inches(1.2), Inches(0.05), Inches(5.8), T.RGB_ACCENT)

    # Left column
    txL = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.7))
    tfL = txL.text_frame
    tfL.word_wrap = True
    for line in left_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        p = tfL.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = _rgb(T.RGB_TEXT_DARK)
        p.font.name = T.FONT_CN
        p.space_before = Pt(4)
    if tfL.paragraphs[0].text == "":
        tfL.paragraphs[0]._element.getparent().remove(tfL.paragraphs[0]._element)

    # Right column
    txR = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.7))
    tfR = txR.text_frame
    tfR.word_wrap = True
    for line in right_text.split("\n"):
        line = line.strip()
        if not line:
            continue
        p = tfR.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = _rgb(T.RGB_TEXT_DARK)
        p.font.name = T.FONT_CN
        p.space_before = Pt(4)
    if tfR.paragraphs[0].text == "":
        tfR.paragraphs[0]._element.getparent().remove(tfR.paragraphs[0]._element)

# ---------------------------------------------------------------------------
# Chart data table slide (ECharts option → table)
# ---------------------------------------------------------------------------

def build_chart_table_slide(prs: Presentation, chart_option: dict[str, Any]) -> None:
    """Extract ECharts option data and present as a table slide."""
    x_axis = chart_option.get("xAxis", {})
    categories = x_axis.get("data", []) if isinstance(x_axis, dict) else []
    series_list = chart_option.get("series", [])
    if not categories or not series_list:
        return

    title_obj = chart_option.get("title", {})
    chart_title = title_obj.get("text", "图表数据") if isinstance(title_obj, dict) else "图表数据"

    slide = _blank(prs)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 chart_title, font_size=T.SIZE_H2, bold=True, color=T.RGB_PRIMARY)

    headers = ["类别"] + [s.get("name", f"系列{i+1}") for i, s in enumerate(series_list)]
    n_rows = min(len(categories), 15)
    n_cols = len(headers)

    tbl_shape = slide.shapes.add_table(1 + n_rows, n_cols,
                                       Inches(0.5), Inches(1.2),
                                       Inches(9), Inches(min(n_rows * 0.45 + 0.5, 5.8)))
    table = tbl_shape.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(T.RGB_SECONDARY)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(T.SIZE_TABLE_HEADER)
            p.font.bold = True
            p.font.color.rgb = _rgb(T.RGB_WHITE)

    for ri in range(n_rows):
        table.cell(ri + 1, 0).text = str(categories[ri])
        for si, s in enumerate(series_list):
            data_list = s.get("data", [])
            val = data_list[ri] if ri < len(data_list) else None
            table.cell(ri + 1, si + 1).text = _fmt_number(val)
        if ri % 2 == 1:
            for ci in range(n_cols):
                table.cell(ri + 1, ci).fill.solid()
                table.cell(ri + 1, ci).fill.fore_color.rgb = _rgb(T.RGB_BG_LIGHT)
        for ci in range(n_cols):
            for p in table.cell(ri + 1, ci).text_frame.paragraphs:
                p.font.size = Pt(T.SIZE_TABLE_BODY)
                p.font.name = T.FONT_NUM

# ---------------------------------------------------------------------------
# Summary slide
# ---------------------------------------------------------------------------

def build_summary_slide(prs: Presentation, conclusions: list[str]) -> None:
    slide = _blank(prs)
    _set_bg(slide, T.RGB_PRIMARY)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                 "核心结论与建议", font_size=24, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Accent underline
    _add_rect(slide, Inches(3.5), Inches(1.1), Inches(3), Inches(0.06), T.RGB_ACCENT)

    for i, text in enumerate(conclusions[:6]):
        _add_textbox(slide, Inches(1), Inches(1.5 + i * 0.9), Inches(8), Inches(0.8),
                     f"  {text}", font_size=14)

# ---------------------------------------------------------------------------
# Thank-you slide
# ---------------------------------------------------------------------------

def build_thank_you_slide(prs: Presentation) -> None:
    slide = _blank(prs)
    _set_bg(slide, T.RGB_PRIMARY)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.5),
                 "谢谢观看", font_size=T.SIZE_TITLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_rect(slide, Inches(4), Inches(4.0), Inches(2), Inches(0.06), T.RGB_ACCENT)
    _add_textbox(slide, Inches(1), Inches(4.5), Inches(8), Inches(1),
                 "THANK YOU", font_size=18,
                 color=T.RGB_ACCENT, alignment=PP_ALIGN.CENTER,
                 font_name=T.FONT_NUM)
