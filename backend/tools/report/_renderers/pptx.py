"""PPTX BlockRenderer — Step 5 (python-pptx fallback path).

Output is structurally equivalent to the previous
``_build_pptx_deterministic`` path (guarded by Step 0 baseline).

Two design notes that differ from MarkdownRenderer / DocxRenderer:
1. **Section buffer mode**: PPTX legacy decides slide composition
   per-section ("narratives + stats → 两栏 slide"), not per-block. This
   renderer collects all blocks of a section in ``begin_section`` /
   ``emit_*`` then renders the slides in ``end_section``.
2. **PptxGenJS path stays in legacy ``ReportContent``**: the Node
   bridge takes a whole ReportContent at once and is migrated as part
   of Sprint 3 visual work, not Step 5. Here we only own the python-pptx
   path that ``pptx_gen.py`` falls back to.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from backend.tools._field_labels import metric_label
from backend.tools.report import _pptx_slides as S
from backend.tools.report import _theme as T
from backend.tools.report._block_renderer import BlockRendererBase
from backend.tools.report._outline import KPIItem
from backend.tools.report._outline import (
    Asset,
    ChartBlock,
    ChartTablePairBlock,
    ComparisonGridBlock,
    GrowthIndicatorsBlock,
    KpiRowBlock,
    OutlineSection,
    ParagraphBlock,
    ReportOutline,
    SectionCoverBlock,
    StatsAsset,
    TableBlock,
)


def _stats_to_text(summary_stats: dict[str, Any]) -> str:
    lines: list[str] = []
    for col, vals in summary_stats.items():
        if not isinstance(vals, dict):
            continue
        mean = vals.get("mean")
        std = vals.get("std")
        if mean is not None:
            line = f"{col}：{metric_label('mean')} {mean:,.2f}"
            if std is not None:
                line += f"  {metric_label('std')} {std:,.2f}"
            lines.append(line)
    return "\n".join(lines) if lines else "暂无统计数据"


class PptxBlockRenderer(BlockRendererBase):
    _step_label = "Step 5"

    def __init__(self, theme=None) -> None:  # noqa: ANN001
        super().__init__(theme=theme)
        self._prs = Presentation()
        self._prs.slide_width = Inches(T.SLIDE_WIDTH)
        self._prs.slide_height = Inches(T.SLIDE_HEIGHT)

        self._title: str = ""
        self._numbered_count: int = 0
        self._is_appendix: bool = False
        self._current_section_name: str = ""

        # Per-section buffers, reset in begin_section
        self._narratives: list[str] = []
        self._stats: list[dict[str, Any]] = []
        self._growth: list[dict[str, Any]] = []
        self._charts: list[dict[str, Any]] = []

        # Appendix + cross-section state
        self._appendix_paragraphs: list[str] = []
        self._all_narratives: list[str] = []  # for fallback summary text

    # ---- Lifecycle -----------------------------------------------------

    def begin_document(self, outline: ReportOutline) -> None:
        self._title = outline.metadata.get("title", "")
        author = outline.metadata.get("author", "")
        date = outline.metadata.get("date", "")

        S.build_cover_slide(self._prs, self._title, author, date)

        # TOC excludes appendix — equivalent to legacy where ReportContent
        # had no appendix section.
        section_names = [
            s.name for s in outline.sections if s.role != "appendix"
        ]
        S.build_toc_slide(self._prs, section_names)

        if outline.kpi_summary:
            self._add_kpi_overview_slide(outline.kpi_summary)

    def end_document(self) -> bytes:
        buf = io.BytesIO()
        self._prs.save(buf)
        return buf.getvalue()

    def begin_section(self, section: OutlineSection, index: int) -> None:
        if section.role == "appendix":
            self._is_appendix = True
            self._appendix_paragraphs = []
        else:
            self._is_appendix = False
            self._numbered_count += 1
            self._current_section_name = section.name
            # Phase 3.1: divider slide is emitted by ``emit_section_cover``
            # (driven by the SectionCoverBlock that legacy converter / LLM
            # planner inserts as the section's first block). begin_section
            # is now state-only; cover behaviour is data-driven.
            self._narratives = []
            self._stats = []
            self._growth = []
            self._charts = []

    def end_section(self, section: OutlineSection, index: int) -> None:
        if self._is_appendix:
            self._render_summary_and_thanks()
        else:
            self._render_section_combo()

    # ---- Block emitters ------------------------------------------------

    def emit_kpi_row(self, block: KpiRowBlock) -> None:
        # Mid-section KPI rows are not represented in legacy PPTX
        # (only the global kpi_summary above section 1 is).
        return None

    def emit_paragraph(self, block: ParagraphBlock) -> None:
        # Phase 4.1 — prefix callouts with emoji marker so the buffer-mode
        # narrative slide still surfaces the warn/info hierarchy without
        # breaking the existing buffer schema. A standalone callout shape
        # is Phase 5.1 component-library work.
        text = block.text
        if block.style == "callout-warn":
            text = f"⚠ 注意: {text}"
        elif block.style == "callout-info":
            text = f"💡 提示: {text}"
        if self._is_appendix:
            self._appendix_paragraphs.append(text)
        else:
            self._narratives.append(text)
            self._all_narratives.append(text)

    def emit_table(self, block: TableBlock, asset: Asset) -> None:
        if isinstance(asset, StatsAsset):
            self._stats.append(asset.summary_stats)
        # TableAsset (DataFrame) is not rendered by legacy PPTX.

    def emit_chart(self, block: ChartBlock, asset: Asset) -> None:
        option = getattr(asset, "option", None)
        if isinstance(option, dict):
            self._charts.append(option)

    def emit_chart_table_pair(
        self,
        block: ChartTablePairBlock,
        chart_asset: Asset,
        table_asset: Asset,
    ) -> None:
        synth_chart = ChartBlock(
            block_id=block.block_id, asset_id=block.chart_asset_id,
        )
        synth_table = TableBlock(
            block_id=block.block_id, asset_id=block.table_asset_id,
        )
        self.emit_chart(synth_chart, chart_asset)
        self.emit_table(synth_table, table_asset)

    def emit_comparison_grid(self, block: ComparisonGridBlock) -> None:
        """Phase 3.2 — dedicated comparison-grid slide via the
        python-pptx fallback path (Node bridge unavailable)."""
        if not block.columns:
            return
        slide_title = self._current_section_name or "对比分析"
        S.build_comparison_grid_slide(self._prs, slide_title, block.columns)

    def emit_growth_indicators(self, block: GrowthIndicatorsBlock) -> None:
        if block.growth_rates:
            self._growth.append(block.growth_rates)

    def emit_section_cover(self, block: SectionCoverBlock) -> None:
        """Paint the section divider slide using the SectionCoverBlock's
        (index, title). Since begin_section is now state-only this is the
        sole owner of the divider visual."""
        S.build_section_divider_slide(self._prs, block.index, block.title)

    # ---- Helpers -------------------------------------------------------

    def _render_section_combo(self) -> None:
        """Replicate legacy section composition: growth slides → narrative+stats
        combo → chart slides."""
        for gr in self._growth:
            S.build_kpi_cards_slide(self._prs, self._current_section_name, gr)

        if self._narratives and self._stats:
            nar_text = "\n\n".join(self._narratives)
            stats_text = _stats_to_text(self._stats[0])
            S.build_two_column_slide(
                self._prs, self._current_section_name, nar_text, stats_text,
            )
            for st in self._stats:
                S.build_stats_table_slide(
                    self._prs,
                    f"{self._current_section_name} - 统计数据",
                    st,
                )
        elif self._narratives:
            S.build_narrative_slide(
                self._prs,
                self._current_section_name,
                "\n\n".join(self._narratives),
            )
        elif self._stats:
            for st in self._stats:
                S.build_stats_table_slide(
                    self._prs,
                    f"{self._current_section_name} - 统计数据",
                    st,
                )

        for ci in self._charts:
            S.build_chart_table_slide(self._prs, ci)

    def _render_summary_and_thanks(self) -> None:
        """Replicate legacy summary slide logic: prefer summary_items
        (truncated to 120 chars); fall back to first long narrative;
        finally a constant default sentence."""
        conclusions: list[str] = [
            (t[:120] + "...") if len(t) > 120 else t
            for t in self._appendix_paragraphs
        ]
        if not conclusions:
            for nar in self._all_narratives:
                if len(nar) > 20:
                    conclusions.append(nar[:100] + "...")
                    break
        if not conclusions:
            conclusions = ["数据分析完成，详见各章节内容"]

        S.build_summary_slide(self._prs, conclusions)
        S.build_thank_you_slide(self._prs)

    def _add_kpi_overview_slide(self, kpis: list[KPIItem]) -> None:
        """Inline KPI-overview slide builder — replicates the legacy
        block in pptx_gen.py without adding a new slide builder.
        """
        from backend.tools.report._pptx_slides import _add_rect, _add_textbox

        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*T.RGB_BG_LIGHT)

        _add_textbox(
            slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
            "核心经营指标", font_size=22, bold=True,
            color=T.RGB_PRIMARY, alignment=PP_ALIGN.LEFT,
        )

        n = min(len(kpis), 4)
        if n == 0:
            return
        card_w = 8.0 / n
        for i, kpi in enumerate(kpis[:n]):
            cx = 1.0 + i * card_w
            _add_rect(
                slide, Inches(cx), Inches(1.3),
                Inches(card_w - 0.2), Inches(4.5),
                T.RGB_BG_LIGHT,
            )
            _add_textbox(
                slide, Inches(cx + 0.1), Inches(1.5),
                Inches(card_w - 0.4), Inches(0.4),
                kpi.label, font_size=10, color=T.RGB_NEUTRAL,
                alignment=PP_ALIGN.CENTER,
            )
            color = T.RGB_POSITIVE if kpi.trend == "positive" else (
                T.RGB_NEGATIVE if kpi.trend == "negative" else T.RGB_PRIMARY
            )
            _add_textbox(
                slide, Inches(cx + 0.1), Inches(2.0),
                Inches(card_w - 0.4), Inches(1.2),
                kpi.value, font_size=36, bold=True, color=color,
                alignment=PP_ALIGN.CENTER, font_name=T.FONT_NUM,
            )
            if kpi.sub:
                _add_textbox(
                    slide, Inches(cx + 0.1), Inches(3.3),
                    Inches(card_w - 0.4), Inches(0.4),
                    kpi.sub, font_size=9, color=T.RGB_NEUTRAL,
                    alignment=PP_ALIGN.CENTER,
                )
