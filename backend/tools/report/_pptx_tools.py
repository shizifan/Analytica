"""PPTX tool factory — wraps _pptx_slides builder functions as LangChain tools.

Usage inside ``PptxReportSkill.execute()``::

    tools = make_pptx_tools(prs, content)
    success = await run_report_agent(llm, tools, PPTX_SYSTEM_PROMPT, msg)
"""
from __future__ import annotations

from typing import Any

from pptx import Presentation
from langchain_core.tools import tool

from backend.tools.report import _pptx_slides as S
from backend.tools.report._agent_loop import FINALIZE_SENTINEL
from backend.tools.report._content_collector import (
    ChartDataItem,
    GrowthItem,
    NarrativeItem,
    ReportContent,
    StatsTableItem,
)

# ---------------------------------------------------------------------------
# System prompt
# ---------------------------------------------------------------------------

PPTX_SYSTEM_PROMPT = """\
你是一位专业的演示文稿设计师。你的任务是使用提供的工具将分析内容组合成一份视觉精美、结构清晰的 PowerPoint 演示文稿。

## 排版规范
1. 首先添加封面幻灯片（add_cover_slide），然后添加目录页（add_toc_slide）
2. 按章节顺序处理每个 Section：
   a. 先添加章节分隔页（add_section_divider）
   b. 如果有增长率数据，添加 KPI 卡片页（add_kpi_cards_slide）
   c. 如果同时有叙述文本和统计数据，考虑使用双栏布局（add_two_column_slide）
   d. 否则分别添加叙述页（add_narrative_slide）和统计表格页（add_stats_table_slide）
   e. 最后添加图表数据页（add_chart_table_slide）
3. 添加总结页（add_summary_slide）和感谢页（add_thank_you_slide）
4. 完成后必须调用 finalize_presentation

## 设计原则
- 每页一个核心信息，避免信息过载
- KPI 卡片放在叙述文本之前，先给出关键数据再展开分析
- 双栏布局适合"左文右数据"的对比呈现
- 空章节可以跳过
- 不要自行编造数据

## 重要规则
- 通过 section_index 和 item_index 引用内容，对应用户消息中的 Section 编号和 [] 编号
"""


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _stats_to_text(summary_stats: dict[str, Any]) -> str:
    """Convert summary_stats dict to readable text for two-column layout."""
    lines = []
    for col, vals in summary_stats.items():
        if not isinstance(vals, dict):
            continue
        parts = []
        for k in ("mean", "median", "std", "min", "max"):
            v = vals.get(k)
            if v is not None:
                label = {"mean": "均值", "median": "中位数", "std": "标准差", "min": "最小值", "max": "最大值"}[k]
                parts.append(f"{label}: {v:,.2f}" if isinstance(v, float) and abs(v) >= 1 else f"{label}: {v}")
        lines.append(f"{col}\n  " + " | ".join(parts))
    return "\n\n".join(lines)


# ---------------------------------------------------------------------------
# Tool factory
# ---------------------------------------------------------------------------

def make_pptx_tools(prs: Presentation, content: ReportContent) -> list:
    """Return LangChain tools that operate on *prs* via closure."""

    def _lookup(section_index: int, item_index: int, expected_type: type) -> Any:
        if section_index < 0 or section_index >= len(content.sections):
            raise IndexError(f"section_index={section_index} 越界 (0-{len(content.sections)-1})")
        sec = content.sections[section_index]
        if item_index < 0 or item_index >= len(sec.items):
            raise IndexError(f"item_index={item_index} 越界，Section '{sec.name}' 共 {len(sec.items)} 项")
        item = sec.items[item_index]
        if not isinstance(item, expected_type):
            raise TypeError(
                f"Section[{section_index}][{item_index}] 是 {type(item).__name__}，"
                f"期望 {expected_type.__name__}"
            )
        return item

    # ── Tools ──────────────────────────────────────────────────────

    @tool
    def add_cover_slide(title: str, author: str, date: str) -> str:
        """添加演示文稿封面页，包含标题、作者和日期。"""
        S.build_cover_slide(prs, title, author, date)
        return f"✓ 封面页已添加：{title}"

    @tool
    def add_toc_slide() -> str:
        """添加目录页，列出所有章节名称。"""
        section_names = [s.name for s in content.sections]
        S.build_toc_slide(prs, section_names)
        return f"✓ 目录页已添加（{len(section_names)} 个章节）"

    @tool
    def add_section_divider(number: int, title: str) -> str:
        """添加章节分隔页，显示章节编号和标题。"""
        S.build_section_divider_slide(prs, number, title)
        return f"✓ 章节分隔页已添加：{number:02d} {title}"

    @tool
    def add_narrative_slide(section_index: int, item_index: int) -> str:
        """添加叙述分析幻灯片（长文本会自动分页）。"""
        item = _lookup(section_index, item_index, NarrativeItem)
        sec_name = content.sections[section_index].name
        S.build_narrative_slide(prs, sec_name, item.text)
        return f"✓ 叙述幻灯片已添加（{len(item.text)} 字符）"

    @tool
    def add_kpi_cards_slide(section_index: int, item_index: int) -> str:
        """添加 KPI 卡片幻灯片（展示同比/环比增长率）。"""
        item = _lookup(section_index, item_index, GrowthItem)
        sec_name = content.sections[section_index].name
        S.build_kpi_cards_slide(prs, sec_name, item.growth_rates)
        return f"✓ KPI 卡片已添加（{len(item.growth_rates)} 项指标）"

    @tool
    def add_stats_table_slide(section_index: int, item_index: int) -> str:
        """添加统计数据表格幻灯片。"""
        item = _lookup(section_index, item_index, StatsTableItem)
        sec_name = content.sections[section_index].name
        S.build_stats_table_slide(prs, sec_name, item.summary_stats)
        return f"✓ 统计表格幻灯片已添加（{len(item.summary_stats)} 列）"

    @tool
    def add_two_column_slide(title: str, left_text: str, right_text: str) -> str:
        """添加双栏布局幻灯片（左文右数据）。"""
        S.build_two_column_slide(prs, title, left_text, right_text)
        return f"✓ 双栏幻灯片已添加：{title}"

    @tool
    def add_chart_table_slide(section_index: int, item_index: int) -> str:
        """添加图表数据表格幻灯片（从 ECharts 配置提取数据）。"""
        item = _lookup(section_index, item_index, ChartDataItem)
        S.build_chart_table_slide(prs, item.option)
        return f"✓ 图表数据幻灯片已添加：{item.title or '图表'}"

    @tool
    def add_summary_slide() -> str:
        """添加核心结论与建议幻灯片。"""
        conclusions = [si.text[:120] + "..." if len(si.text) > 120 else si.text
                       for si in content.summary_items]
        if not conclusions:
            conclusions = ["数据分析完成，详见各章节内容"]
        S.build_summary_slide(prs, conclusions)
        return f"✓ 总结幻灯片已添加（{len(conclusions)} 条结论）"

    @tool
    def add_thank_you_slide() -> str:
        """添加感谢观看幻灯片。"""
        S.build_thank_you_slide(prs)
        return "✓ 感谢页已添加"

    @tool
    def finalize_presentation() -> str:
        """完成演示文稿编排，结束工具调用。必须在所有内容添加完毕后调用。"""
        return FINALIZE_SENTINEL

    return [
        add_cover_slide,
        add_toc_slide,
        add_section_divider,
        add_narrative_slide,
        add_kpi_cards_slide,
        add_stats_table_slide,
        add_two_column_slide,
        add_chart_table_slide,
        add_summary_slide,
        add_thank_you_slide,
        finalize_presentation,
    ]
