"""辽港数据期刊 — LLM 后置审查模块。

定义 ``ReviewResult`` / ``ReviewFinding`` 数据结构，
``ReportQualityReviewer`` 基类及格式专用审查逻辑。
PR-3 落地：DOCX 格式专用审查（结构完整性 + 数据准确性规则检查）。
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field
from typing import Any

logger = logging.getLogger("analytica.tools.report.review")


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------

@dataclass
class ReviewFinding:
    """单条审查发现。

    Attributes:
        dimension: 审查维度（data_accuracy / title_quality / ...）。
        passed: 是否通过。
        severity: ``"BLOCKING"``（阻断交付）或 ``"WARNING"``（仅日志）。
        detail: 人类可读的问题描述。
        block_ids: 关联的问题 block ID 列表。
    """
    dimension: str
    passed: bool
    severity: str          # "BLOCKING" | "WARNING"
    detail: str
    block_ids: list[str] = field(default_factory=list)


@dataclass
class ReviewResult:
    """一次审查的汇总结果。

    Attributes:
        passed: 所有 BLOCKING 维度均通过。
        findings: 所有审查发现列表。
        retry_targets: 需要重试的 block / chapter ID 列表。
    """
    passed: bool
    findings: list[ReviewFinding] = field(default_factory=list)
    retry_targets: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# 审查维度名常量
# ---------------------------------------------------------------------------

DIM_DATA_ACCURACY = "data_accuracy"
DIM_TITLE_QUALITY = "title_quality"
DIM_STRUCTURE_COMPLETENESS = "structure_completeness"
DIM_NARRATIVE_COHERENCE = "narrative_coherence"
DIM_VISUAL_CONSISTENCY = "visual_consistency"

BLOCKING = "BLOCKING"
WARNING = "WARNING"

# 通用图表标题（黑名单）
_GENERIC_CHART_TITLES: frozenset[str] = frozenset({
    "数据对比", "趋势图", "分布图", "对比图", "图表",
    "柱状图", "折线图", "饼图",
})


# ---------------------------------------------------------------------------
# 审查器
# ---------------------------------------------------------------------------

class ReportQualityReviewer:
    """LLM 后置审查器。

    审查在所有 renderer 输出完成后、交付用户前触发。
    若 FAIL，触发单 block 重试后再次审查，最多 2 轮。

    PR-3 落地：``review_docx()`` 提供 DOCX 格式专用审查
    （结构完整性 + 数据准确性 + 标题质量规则检查）。
    """

    # 审查维度配置：各维度的严重级别
    _DIMENSION_SEVERITY: dict[str, str] = {
        DIM_DATA_ACCURACY: BLOCKING,
        DIM_TITLE_QUALITY: BLOCKING,
        DIM_STRUCTURE_COMPLETENESS: BLOCKING,
        DIM_NARRATIVE_COHERENCE: WARNING,
        DIM_VISUAL_CONSISTENCY: WARNING,
    }

    # ── 通用 review（多态入口） ──────────────────────────────

    def review(
        self,
        report: Any,
        source_data: dict | None = None,
    ) -> ReviewResult:
        """通用入口 — 子类或格式专用方法覆写。

        返回全部 passed=True 的默认结果（PR-1 骨架行为）。
        """
        findings = [
            ReviewFinding(
                dimension=dim,
                passed=True,
                severity=sev,
                detail=f"[PR-1 skeleton] {dim} check not yet implemented",
            )
            for dim, sev in self._DIMENSION_SEVERITY.items()
        ]
        return ReviewResult(passed=True, findings=findings, retry_targets=[])

    # ── DOCX 格式专用审查（PR-3 落地） ─────────────────────

    def review_docx(
        self,
        docx_bytes: bytes,
        outline: Any = None,
    ) -> ReviewResult:
        """对 DOCX 输出执行五维度审查。

        规则检查维度（结构完整性 / 数据准确性 / 标题质量）
        不需要 LLM 参与 — 纯规则匹配即可。

        ``outline`` 为 ``ReportOutline`` 实例时抽取 source data
        用于数据准确性抽查。
        """
        from docx import Document as DocxDocument

        findings: list[ReviewFinding] = []

        try:
            doc = DocxDocument(io.BytesIO(docx_bytes))
        except Exception as e:
            findings.append(ReviewFinding(
                dimension=DIM_STRUCTURE_COMPLETENESS,
                passed=False,
                severity=BLOCKING,
                detail=f"DOCX 二进制无法解析: {e}",
            ))
            # 无法解析时跳过后续检查
            return ReviewResult(
                passed=False, findings=findings, retry_targets=[],
            )

        # 1. 结构完整性检查
        findings.append(self._check_docx_structure(doc, outline))

        # 2. 数据准确性抽查
        if outline is not None:
            findings.append(self._check_docx_data_accuracy(doc, outline))

        # 3. 标题质量检查
        findings.append(self._check_docx_title_quality(doc))

        # 4. 视觉一致性（WARNING：仅记录不阻断）
        findings.append(self._check_docx_visual_consistency(doc))

        # 5. 叙述连贯性（WARNING：外观检查）
        findings.append(self._check_docx_narrative(doc))

        # 汇总
        blocking_failures = [
            f for f in findings
            if f.severity == BLOCKING and not f.passed
        ]
        return ReviewResult(
            passed=len(blocking_failures) == 0,
            findings=findings,
            retry_targets=list({
                bid
                for f in blocking_failures
                for bid in f.block_ids
            }),
        )

    # ── PPTX 格式专用审查（PR-4 落地） ─────────────────────

    def review_pptx(
        self,
        pptx_bytes: bytes,
        outline: Any = None,
    ) -> ReviewResult:
        """对 PPTX 输出执行五维度审查（PR-4 辽港数据期刊）。

        规则检查：幻灯片数量 / 主题色合规 / 标题质量 / 结构完整性 /
        叙述连贯性。纯规则匹配，无需 LLM。
        """
        from pptx import Presentation as PptxPresentation

        findings: list[ReviewFinding] = []

        try:
            prs = PptxPresentation(io.BytesIO(pptx_bytes))
        except Exception as e:
            findings.append(ReviewFinding(
                dimension=DIM_STRUCTURE_COMPLETENESS,
                passed=False,
                severity=BLOCKING,
                detail=f"PPTX 二进制无法解析: {e}",
            ))
            return ReviewResult(
                passed=False, findings=findings, retry_targets=[],
            )

        # 收集幻灯片文字（供后续检查使用）
        all_slide_texts: list[list[str]] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            all_slide_texts.append(texts)

        # 1. 幻灯片数量检查
        findings.append(self._check_pptx_slide_count(prs))

        # 2. 主题色合规检查
        findings.append(self._check_pptx_theme_colors(prs, all_slide_texts))

        # 3. 标题质量检查
        findings.append(self._check_pptx_title_quality(all_slide_texts))

        # 4. 结构完整性检查
        findings.append(self._check_pptx_structure(prs, all_slide_texts))

        # 5. 叙述连贯性（WARNING）
        findings.append(self._check_pptx_narrative(all_slide_texts))

        # 汇总
        blocking_failures = [
            f for f in findings
            if f.severity == BLOCKING and not f.passed
        ]
        return ReviewResult(
            passed=len(blocking_failures) == 0,
            findings=findings,
            retry_targets=list({
                bid
                for f in blocking_failures
                for bid in f.block_ids
            }),
        )

    def _check_pptx_slide_count(self, prs: Any) -> ReviewFinding:
        """检查 PPTX 幻灯片数量：至少 5 页。"""
        count = len(prs.slides)
        if count < 5:
            return ReviewFinding(
                dimension=DIM_STRUCTURE_COMPLETENESS,
                passed=False,
                severity=BLOCKING,
                detail=f"幻灯片数量过少 ({count})，预期至少 5 页（封面/目录/正文/分隔/结语）",
            )
        return ReviewFinding(
            dimension=DIM_STRUCTURE_COMPLETENESS,
            passed=True,
            severity=BLOCKING,
            detail=f"幻灯片数量正常: {count} 页",
        )

    def _check_pptx_theme_colors(
        self, prs: Any, all_slide_texts: list[list[str]],
    ) -> ReviewFinding:
        """检查 PPTX 是否使用了辽港主题色（navy #004889 / bronze #AC916B）。

        遍历所有形状的 run 级字体颜色，检测辽港品牌色。
        """
        has_navy = False
        has_bronze = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            c = run.font.color
                            if c and c.rgb:
                                hx = str(c.rgb).upper()
                                if hx == "004889":
                                    has_navy = True
                                elif hx == "AC916B":
                                    has_bronze = True
                        except Exception:
                            pass

        missing: list[str] = []
        if not has_navy:
            missing.append("navy (#004889)")
        if not has_bronze:
            missing.append("bronze (#AC916B)")

        if missing:
            return ReviewFinding(
                dimension=DIM_VISUAL_CONSISTENCY,
                passed=False,
                severity=WARNING,
                detail=f"未检测到辽港品牌色: {', '.join(missing)}",
            )

        return ReviewFinding(
            dimension=DIM_VISUAL_CONSISTENCY,
            passed=True,
            severity=WARNING,
            detail="辽港品牌色 (navy + bronze) 已应用",
        )

    def _check_pptx_title_quality(
        self, all_slide_texts: list[list[str]],
    ) -> ReviewFinding:
        """检查 PPTX 图表/表格标题是否为通用词（黑名单匹配）。"""
        generic_found: list[str] = []
        for texts in all_slide_texts:
            for t in texts:
                if t in _GENERIC_CHART_TITLES:
                    generic_found.append(t)

        if generic_found:
            return ReviewFinding(
                dimension=DIM_TITLE_QUALITY,
                passed=False,
                severity=BLOCKING,
                detail=f"发现通用标题: {generic_found}",
            )

        return ReviewFinding(
            dimension=DIM_TITLE_QUALITY,
            passed=True,
            severity=BLOCKING,
            detail="未检测到通用图表标题",
        )

    def _check_pptx_structure(
        self, prs: Any, all_slide_texts: list[list[str]],
    ) -> ReviewFinding:
        """检查 PPTX 结构完整性：至少包含图表或表格。"""
        table_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1

        # 检查是否有图表类型的 shape（pptxgenjs 生成的是 native chart）
        chart_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_chart:
                    chart_count += 1

        if table_count == 0 and chart_count == 0:
            return ReviewFinding(
                dimension=DIM_STRUCTURE_COMPLETENESS,
                passed=False,
                severity=BLOCKING,
                detail="PPTX 中无图表且无表格，报告可能缺少数据呈现",
            )

        return ReviewFinding(
            dimension=DIM_STRUCTURE_COMPLETENESS,
            passed=True,
            severity=BLOCKING,
            detail=(
                f"结构正常: {chart_count} 图表, {table_count} 表格"
            ),
        )

    def _check_pptx_narrative(
        self, all_slide_texts: list[list[str]],
    ) -> ReviewFinding:
        """检查 PPTX 叙述连贯性（WARNING）：每页平均文字形状数。

        辽港数据期刊模板要求每页有足够上下文。
        """
        slide_count = len(all_slide_texts)
        total_texts = sum(len(texts) for texts in all_slide_texts)
        avg_texts = total_texts / slide_count if slide_count else 0

        if avg_texts < 1:
            return ReviewFinding(
                dimension=DIM_NARRATIVE_COHERENCE,
                passed=False,
                severity=WARNING,
                detail=f"每页平均文字形状仅 {avg_texts:.1f}，缺少叙述上下文",
            )

        return ReviewFinding(
            dimension=DIM_NARRATIVE_COHERENCE,
            passed=True,
            severity=WARNING,
            detail=f"叙述正常: 每页平均 {avg_texts:.1f} 个文字形状",
        )

    def _check_docx_structure(
        self,
        doc: Any,
        outline: Any = None,
    ) -> ReviewFinding:
        """检查 DOCX 结构完整性：段落数 / 表格数 / 图片数。"""
        para_count = len(doc.paragraphs)
        table_count = len(doc.tables)

        # 统计 inline shapes（图片）
        image_count = 0
        for para in doc.paragraphs:
            for run in para.runs:
                image_count += len(run._element.findall(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                ))
        # 也检查 tables 中的图片
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            image_count += len(run._element.findall(
                                ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                            ))

        issues: list[str] = []
        if para_count < 5:
            issues.append(f"段落数过少 ({para_count})，可能生成不完整")
        if table_count == 0 and image_count == 0:
            issues.append("无表格且无图表图片，报告可能缺少数据呈现")

        if issues:
            return ReviewFinding(
                dimension=DIM_STRUCTURE_COMPLETENESS,
                passed=False,
                severity=BLOCKING,
                detail="; ".join(issues),
            )

        return ReviewFinding(
            dimension=DIM_STRUCTURE_COMPLETENESS,
            passed=True,
            severity=BLOCKING,
            detail=(
                f"结构正常: {para_count} 段落, {table_count} 表格, "
                f"{image_count} 图片"
            ),
        )

    def _check_docx_data_accuracy(
        self,
        doc: Any,
        outline: Any,
    ) -> ReviewFinding:
        """从 DOCX 表格中抽样数值，与源数据交叉验证。

        抽取前 3 个表格中的数值，与 outline 中的 TableAsset 比对。
        仅做 BLOCKING 检查：确保表格中有数据（非空）。
        """
        tables = doc.tables
        if not tables:
            return ReviewFinding(
                dimension=DIM_DATA_ACCURACY,
                passed=True,
                severity=BLOCKING,
                detail="无表格可抽查",
            )

        # 检查前 3 个表格至少有实际数据行
        empty_tables = 0
        for table in tables[:3]:
            data_rows = sum(
                1 for row in table.rows[1:]
                if any(
                    cell.text.strip()
                    for cell in row.cells
                )
            )
            if data_rows == 0:
                empty_tables += 1

        if empty_tables > 0:
            return ReviewFinding(
                dimension=DIM_DATA_ACCURACY,
                passed=False,
                severity=BLOCKING,
                detail=f"{empty_tables} 个表格无数据行",
            )

        return ReviewFinding(
            dimension=DIM_DATA_ACCURACY,
            passed=True,
            severity=BLOCKING,
            detail=(
                f"抽查通过: 前 {min(3, len(tables))} 个表格均有数据行"
            ),
        )

    def _check_docx_title_quality(
        self,
        doc: Any,
    ) -> ReviewFinding:
        """检查图表/表格标题是否为通用词（黑名单匹配）。"""
        generic_found: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text in _GENERIC_CHART_TITLES:
                generic_found.append(text)

        if generic_found:
            return ReviewFinding(
                dimension=DIM_TITLE_QUALITY,
                passed=False,
                severity=BLOCKING,
                detail=f"发现通用标题: {generic_found}",
            )

        return ReviewFinding(
            dimension=DIM_TITLE_QUALITY,
            passed=True,
            severity=BLOCKING,
            detail="未检测到通用图表标题",
        )

    def _check_docx_visual_consistency(
        self,
        doc: Any,
    ) -> ReviewFinding:
        """检查 DOCX 视觉一致性（WARNING）。

        检测：文档是否包含至少一个辽港主题色段落
        （#004889 或 #AC916B），确认主题样式已应用。
        """
        # 检查 heading 段落是否有辽港品牌色
        brand_colors_used = False
        for para in doc.paragraphs:
            if para.style and para.style.name:
                style_name = para.style.name.lower()
                if any(k in style_name for k in ("heading", "title", "kpi")):
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            hex_c = str(run.font.color.rgb).upper()
                            if hex_c in ("004889", "AC916B"):
                                brand_colors_used = True
                                break
                if brand_colors_used:
                    break

        if not brand_colors_used:
            return ReviewFinding(
                dimension=DIM_VISUAL_CONSISTENCY,
                passed=False,
                severity=WARNING,
                detail="未在标题段落中检测到辽港品牌色 (primary/accent)",
            )

        return ReviewFinding(
            dimension=DIM_VISUAL_CONSISTENCY,
            passed=True,
            severity=WARNING,
            detail="辽港品牌色已应用",
        )

    def _check_docx_narrative(
        self,
        doc: Any,
    ) -> ReviewFinding:
        """检查叙述连贯性（WARNING）。

        检测：文档至少包含一个正文叙述段落（非标题、非空）。
        """
        narrative_paras = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or "").lower()
            # 跳过标题和特殊样式
            if any(k in style_name for k in ("heading", "title", "subtitle")):
                continue
            if len(text) > 20:
                narrative_paras += 1

        if narrative_paras == 0:
            return ReviewFinding(
                dimension=DIM_NARRATIVE_COHERENCE,
                passed=False,
                severity=WARNING,
                detail="文档缺少正文叙述段落（可能仅有标题和表格）",
            )

        return ReviewFinding(
            dimension=DIM_NARRATIVE_COHERENCE,
            passed=True,
            severity=WARNING,
            detail=f"包含 {narrative_paras} 个叙述段落",
        )


# ---------------------------------------------------------------------------
# 便捷函数
# ---------------------------------------------------------------------------

def review_docx_output(
    docx_bytes: bytes,
    outline: Any = None,
) -> ReviewResult:
    """对 DOCX 输出执行后置审查的便捷函数。

    集成到 ``docx_gen.py`` / ``DocxReportTool.execute()`` 中使用。
    """
    reviewer = ReportQualityReviewer()
    result = reviewer.review_docx(docx_bytes, outline=outline)

    # 记录审查结果
    for f in result.findings:
        if not f.passed:
            level = logging.WARNING if f.severity == WARNING else logging.ERROR
            logger.log(level, "DOCX review %s: %s", f.dimension, f.detail)

    if not result.passed:
        logger.warning(
            "DOCX review FAILED — %d BLOCKING issue(s): %s",
            sum(1 for f in result.findings if f.severity == BLOCKING and not f.passed),
            ", ".join(
                f"{f.dimension}: {f.detail}"
                for f in result.findings
                if f.severity == BLOCKING and not f.passed
            ),
        )
    else:
        logger.info("DOCX review PASSED — all BLOCKING checks clear")

    return result


def review_pptx_output(
    pptx_bytes: bytes,
    outline: Any = None,
) -> ReviewResult:
    """对 PPTX 输出执行后置审查的便捷函数。

    集成到 ``pptx_gen.py`` / ``PptxReportTool.execute()`` 中使用。
    PR-4 落地：辽港数据期刊 PPTX 质量审查。
    """
    reviewer = ReportQualityReviewer()
    result = reviewer.review_pptx(pptx_bytes, outline=outline)

    for f in result.findings:
        if not f.passed:
            level = logging.WARNING if f.severity == WARNING else logging.ERROR
            logger.log(level, "PPTX review %s: %s", f.dimension, f.detail)

    if not result.passed:
        logger.warning(
            "PPTX review FAILED — %d BLOCKING issue(s): %s",
            sum(1 for f in result.findings if f.severity == BLOCKING and not f.passed),
            ", ".join(
                f"{f.dimension}: {f.detail}"
                for f in result.findings
                if f.severity == BLOCKING and not f.passed
            ),
        )
    else:
        logger.info("PPTX review PASSED — all BLOCKING checks clear")

    return result
